"""PPTX 生成工具 —— 将结构化 Markdown 转换为 PowerPoint 演示文稿。"""

import re
from pathlib import Path

from graph_agent.tools.base import tool

MAX_FILE_SIZE = 1024 * 1024  # 1MB

# ── 样式常量 ────────────────────────────────────────────────
DEEP_BLUE = 0x1A3A5C
DARK_GRAY = 0x333333
MEDIUM_GRAY = 0x666666
LIGHT_GRAY = 0x999999
WHITE = 0xFFFFFF
BLACK = 0x000000
FONT_NAME = "Microsoft YaHei"  # 中英文混排友好


def _resolve_path(path: str) -> Path:
    resolved = Path(path).resolve()
    if ".." in resolved.parts:
        raise ValueError("路径包含非法的上级目录引用")
    return resolved


def _measure_indent(line: str) -> int:
    count = 0
    for ch in line:
        if ch == " ":
            count += 1
        elif ch == "\t":
            count += 2
        else:
            break
    return count // 2


def _parse_markdown(content: str) -> list[dict]:
    """将 Markdown 文本解析为幻灯片数据结构。"""
    lines = content.strip().split("\n")
    slides = []
    current = None
    title_created = False
    just_title = False

    for line in lines:
        stripped = line.strip()

        if not stripped:
            continue

        # 显式分页符
        if stripped in ("---", "***", "___"):
            if current:
                slides.append(current)
                current = None
            just_title = False
            continue

        # 封面标题（仅第一个 #）
        if stripped.startswith("# ") and not title_created:
            if current:
                slides.append(current)
            current = {
                "type": "title",
                "title": stripped[2:].strip(),
                "subtitle": "",
                "items": [],
            }
            title_created = True
            just_title = True
            continue

        # 紧接封面的 ## → 副标题
        if stripped.startswith("## ") and just_title and current and current["type"] == "title":
            current["subtitle"] = stripped[3:].strip()
            just_title = False
            continue

        just_title = False

        # 章节页
        if stripped.startswith("## "):
            if current:
                slides.append(current)
            current = {"type": "section", "title": stripped[3:].strip(), "items": []}
            continue

        # 内容页
        if stripped.startswith("### "):
            if current:
                slides.append(current)
            current = {"type": "content", "title": stripped[4:].strip(), "items": []}
            continue

        # 没有当前幻灯片时自动创建空白内容页
        if current is None:
            current = {"type": "content", "title": "", "items": []}

        # 无序列表
        if re.match(r"^[-*]\s+", stripped):
            text = re.sub(r"^[-*]\s+", "", stripped)
            current["items"].append({
                "type": "bullet", "text": text, "level": _measure_indent(line),
            })
            continue

        # 有序列表
        num_match = re.match(r"^(\d+)[.)]\s+(.+)$", stripped)
        if num_match:
            current["items"].append({"type": "numbered", "text": num_match.group(2).strip()})
            continue

        # 引用
        if stripped.startswith("> "):
            current["items"].append({"type": "quote", "text": stripped[2:].strip()})
            continue

        # 图片
        img_match = re.match(r"^!\[([^\]]*)\]\(([^)]+)\)$", stripped)
        if img_match:
            current["items"].append({
                "type": "image", "alt": img_match.group(1), "path": img_match.group(2),
            })
            continue

        # 普通段落
        current["items"].append({"type": "paragraph", "text": stripped})

    if current:
        slides.append(current)

    return slides


def _render_pptx(slides: list[dict], output_path: str) -> int:
    """使用 python-pptx 将幻灯片数据渲染为 PPTX 文件，返回总页数。"""
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

    prs = Presentation()
    # 设置 16:9 宽屏
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for slide_data in slides:
        slide_type = slide_data["type"]

        if slide_type == "title":
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            _set_title_shape(slide.shapes.title, slide_data["title"], Pt(44))
            if slide_data.get("subtitle"):
                subtitle_shape = _get_subtitle_placeholder(slide)
                if subtitle_shape:
                    _set_text_shape(subtitle_shape, slide_data["subtitle"], Pt(24), MEDIUM_GRAY, PP_ALIGN.CENTER)

        elif slide_type == "section":
            slide = prs.slides.add_slide(prs.slide_layouts[2])
            if slide.shapes.title:
                _set_title_shape(slide.shapes.title, slide_data["title"], Pt(36))
            else:
                # 某些模板的 Section Header 没有 title placeholder，手动添加
                _add_centered_text(slide, slide_data["title"], Pt(36))

        elif slide_type == "content":
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            if slide.shapes.title:
                _set_title_shape(slide.shapes.title, slide_data["title"], Pt(32))

            # 将内容写入正文占位符
            body_shape = _get_body_placeholder(slide)
            if body_shape is not None:
                _fill_content(body_shape, slide_data.get("items", []))
            else:
                # 无正文占位符时手动添加文本框
                _add_content_textbox(slide, slide_data.get("items", []))

    prs.save(str(output_path))
    return len(slides)


# ── 幻灯片渲染辅助函数 ────────────────────────────────────


def _set_title_shape(shape, text: str, font_size):
    """设置标题占位符的文本和样式。"""
    from pptx.util import Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    shape.text = ""
    p = shape.text_frame.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.bold = True
    p.font.color.rgb = RGBColor(
        (DEEP_BLUE >> 16) & 0xFF, (DEEP_BLUE >> 8) & 0xFF, DEEP_BLUE & 0xFF,
    )
    p.font.name = FONT_NAME
    p.alignment = PP_ALIGN.LEFT


def _set_text_shape(shape, text: str, font_size, color: int, alignment):
    """设置文本占位符的内容和样式。"""
    from pptx.util import Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    shape.text = ""
    p = shape.text_frame.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = RGBColor((color >> 16) & 0xFF, (color >> 8) & 0xFF, color & 0xFF)
    p.font.name = FONT_NAME
    p.alignment = alignment


def _get_subtitle_placeholder(slide):
    """获取副标题占位符（idx 1 或 2）。"""
    try:
        return slide.placeholders[1]
    except (KeyError, IndexError):
        pass
    try:
        return slide.placeholders[2]
    except (KeyError, IndexError):
        return None


def _get_body_placeholder(slide):
    """获取正文占位符（通常 idx 1）。"""
    try:
        return slide.placeholders[1]
    except (KeyError, IndexError):
        pass
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 1:
            return shape
    return None


def _add_centered_text(slide, text: str, font_size):
    """在幻灯片中央手动添加居中的文本框。"""
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

    left = Inches(1.5)
    top = Inches(2.5)
    width = Inches(10.333)
    height = Inches(2.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.text_frame.word_wrap = True
    tf = txBox.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.bold = True
    p.font.color.rgb = RGBColor(
        (DEEP_BLUE >> 16) & 0xFF, (DEEP_BLUE >> 8) & 0xFF, DEEP_BLUE & 0xFF,
    )
    p.font.name = FONT_NAME
    tf.paragraphs[0].space_after = Pt(12)


def _add_content_textbox(slide, items: list):
    """当没有正文占位符时手动添加内容文本框。"""
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    left = Inches(1.0)
    top = Inches(1.8)
    width = Inches(11.333)
    height = Inches(5.2)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.text_frame.word_wrap = True
    _fill_content(txBox, items)


def _fill_content(shape, items: list):
    """将内容项列表写入正文占位符的 text_frame。"""
    from pptx.util import Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    tf = shape.text_frame

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
            p.text = ""
        else:
            p = tf.add_paragraph()
        item_type = item["type"]

        if item_type == "bullet":
            p.text = item["text"]
            p.level = min(item.get("level", 0), 2)
            p.font.size = _bullet_font_size(p.level)
            p.font.color.rgb = _bullet_color(p.level)
            p.font.name = FONT_NAME
            p.space_after = Pt(6)

        elif item_type == "numbered":
            p.text = item["text"]
            p.font.size = Pt(20)
            p.font.color.rgb = RGBColor(
                (DARK_GRAY >> 16) & 0xFF, (DARK_GRAY >> 8) & 0xFF, DARK_GRAY & 0xFF,
            )
            p.font.name = FONT_NAME
            p.space_after = Pt(6)

        elif item_type == "paragraph":
            p.text = item["text"]
            p.font.size = Pt(16)
            p.font.color.rgb = RGBColor(
                (DARK_GRAY >> 16) & 0xFF, (DARK_GRAY >> 8) & 0xFF, DARK_GRAY & 0xFF,
            )
            p.font.name = FONT_NAME
            p.space_after = Pt(10)

        elif item_type == "quote":
            p.text = item["text"]
            p.font.size = Pt(14)
            p.font.italic = True
            p.font.color.rgb = RGBColor(
                (LIGHT_GRAY >> 16) & 0xFF, (LIGHT_GRAY >> 8) & 0xFF, LIGHT_GRAY & 0xFF,
            )
            p.font.name = FONT_NAME
            p.space_after = Pt(8)

        elif item_type == "image":
            p.text = f"[图片: {item.get('alt', item['path'])}]"
            p.font.size = Pt(12)
            p.font.italic = True
            p.font.color.rgb = RGBColor(
                (LIGHT_GRAY >> 16) & 0xFF, (LIGHT_GRAY >> 8) & 0xFF, LIGHT_GRAY & 0xFF,
            )
            p.font.name = FONT_NAME
            p.space_after = Pt(4)


def _bullet_font_size(level: int):
    from pptx.util import Pt
    return [Pt(20), Pt(18), Pt(16)][min(level, 2)]


def _bullet_color(level: int):
    from pptx.dml.color import RGBColor
    colors = [
        RGBColor((DARK_GRAY >> 16) & 0xFF, (DARK_GRAY >> 8) & 0xFF, DARK_GRAY & 0xFF),
        RGBColor((MEDIUM_GRAY >> 16) & 0xFF, (MEDIUM_GRAY >> 8) & 0xFF, MEDIUM_GRAY & 0xFF),
        RGBColor((LIGHT_GRAY >> 16) & 0xFF, (LIGHT_GRAY >> 8) & 0xFF, LIGHT_GRAY & 0xFF),
    ]
    return colors[min(level, 2)]


# ── 工具入口 ────────────────────────────────────────────────


@tool(
    "generate_pptx",
    "将结构化Markdown转换为PPTX演示文稿。"
    "参数 markdown_path: Markdown文件路径（与content二选一）, "
    "content: Markdown文本内容（与markdown_path二选一，直接传入内容无需先写文件）, "
    "output_path: 输出PPTX路径（推荐明确指定，未指定时从markdown_path推断）",
    risk_level="low",
)
def generate_pptx(markdown_path: str = "", output_path: str = "",
                  content: str = "") -> str:
    try:
        from pptx import Presentation
    except ImportError:
        return "错误: 缺少依赖库 python-pptx，请执行 pip install python-pptx"

    if not markdown_path and not content:
        return "错误: 必须提供 markdown_path 或 content 参数"

    if content:
        if not output_path:
            return "错误: 使用 content 参数时必须提供 output_path"
        # 将 content 写入临时 Markdown 文件
        out_pptx = Path(output_path).resolve()
        tmp_md = out_pptx.with_suffix(".md")
        try:
            tmp_md.write_text(content, encoding="utf-8")
        except Exception as e:
            return f"错误: 无法写入临时Markdown文件 - {e}"
        path = tmp_md
    else:
        path = _resolve_path(markdown_path)
        if not path.exists():
            return f"错误: Markdown文件不存在 - {markdown_path}"
        if not path.is_file():
            return f"错误: 路径不是文件 - {markdown_path}"
        if path.stat().st_size > MAX_FILE_SIZE:
            return f"错误: Markdown文件超过1MB大小限制"

    try:
        md_content = path.read_text(encoding="utf-8")
    except UnicodeDecodeError:
        return f"错误: 无法以UTF-8编码读取文件 - {path}"
    except Exception as e:
        return f"错误: 读取文件失败 - {e}"

    # 确定输出路径
    if output_path:
        out = Path(output_path).resolve()
    else:
        out = path.with_suffix(".pptx")

    # 确保输出目录存在
    try:
        out.parent.mkdir(parents=True, exist_ok=True)
    except Exception as e:
        return f"错误: 无法创建输出目录 - {e}"

    # 解析 Markdown
    try:
        slides = _parse_markdown(md_content)
    except Exception as e:
        return f"错误: Markdown文件解析失败 - {e}"

    if not slides:
        return "错误: Markdown文件解析后无有效幻灯片内容"

    # 渲染并保存
    try:
        page_count = _render_pptx(slides, str(out))
    except Exception as e:
        return f"错误: PPTX渲染失败 - {e}"

    # 验证文件确实写入磁盘
    if not out.exists():
        return f"错误: PPTX文件保存后验证失败，文件未出现在预期位置: {out}"
    file_size = out.stat().st_size
    if file_size == 0:
        return f"错误: PPTX文件已创建但大小为0字节: {out}"

    # 统计信息
    title_count = sum(1 for s in slides if s["type"] == "title")
    section_count = sum(1 for s in slides if s["type"] == "section")
    content_count = sum(1 for s in slides if s["type"] == "content")

    return (
        f"PPTX生成成功\n"
        f"输出文件: {out}\n"
        f"幻灯片总数: {page_count} 页\n"
        f"  封面页: {title_count} 页\n"
        f"  章节页: {section_count} 页\n"
        f"  内容页: {content_count} 页"
    )
